"""Generate mbo_utilities cheat sheet as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Colors matching the dark theme from GUI screenshots
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)  # Dark blue-gray
ACCENT_YELLOW = RGBColor(0xff, 0xd7, 0x00)  # Gold/yellow
ACCENT_CYAN = RGBColor(0x00, 0xd4, 0xff)  # Cyan
ACCENT_GREEN = RGBColor(0x4e, 0xc9, 0x4e)  # Green
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xcc)
CODE_BG = RGBColor(0x2d, 0x2d, 0x3a)  # Slightly lighter for code blocks


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_YELLOW
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title, content_items, code_examples=None):
    """Add a content slide with bullet points and optional code."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_YELLOW

    # Content area
    y_pos = 0.9

    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.3), Inches(y_pos), Inches(4.5), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if isinstance(item, tuple):
                # (header, description) format
                p.text = f"• {item[0]}"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = ACCENT_CYAN

                p2 = tf.add_paragraph()
                p2.text = f"   {item[1]}"
                p2.font.size = Pt(12)
                p2.font.color.rgb = LIGHT_GRAY
            else:
                p.text = f"• {item}"
                p.font.size = Pt(13)
                p.font.color.rgb = WHITE

    # Code examples on right side
    if code_examples:
        code_box = slide.shapes.add_textbox(Inches(5.0), Inches(y_pos), Inches(4.7), Inches(4.5))
        tf = code_box.text_frame
        tf.word_wrap = True

        for i, (label, code) in enumerate(code_examples):
            if i > 0:
                p = tf.add_paragraph()
                p.text = ""
                p.font.size = Pt(6)

            p = tf.add_paragraph() if i > 0 or tf.paragraphs[0].text else tf.paragraphs[0]
            p.text = f"# {label}"
            p.font.size = Pt(11)
            p.font.color.rgb = ACCENT_GREEN
            p.font.name = "Consolas"

            for line in code.strip().split('\n'):
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(10)
                p.font.color.rgb = WHITE
                p.font.name = "Consolas"

    return slide


def add_cli_slide(prs):
    """Add CLI commands slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CLI Commands"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_YELLOW

    # Commands table-like layout
    commands = [
        ("mbo view [PATH]", "Launch GUI viewer", "--roi 0,1  --widget  --metadata"),
        ("mbo convert IN OUT", "Convert between formats", "-e .zarr  -p 0,1,2  --fix-phase  --register-z"),
        ("mbo info PATH", "Show file metadata", "--metadata"),
        ("mbo scanphase [PATH]", "Analyze scan phase", "-o output/  --format png  --show"),
        ("mbo formats", "List supported formats", ""),
        ("mbo --download-notebook", "Get user guide notebook", "[PATH]"),
    ]

    y = 0.85
    for cmd, desc, opts in commands:
        # Command
        box = slide.shapes.add_textbox(Inches(0.3), Inches(y), Inches(3.2), Inches(0.35))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = cmd
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = "Consolas"

        # Description
        box = slide.shapes.add_textbox(Inches(3.5), Inches(y), Inches(2.5), Inches(0.35))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

        # Options
        if opts:
            box = slide.shapes.add_textbox(Inches(6.0), Inches(y), Inches(3.7), Inches(0.35))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = opts
            p.font.size = Pt(10)
            p.font.color.rgb = LIGHT_GRAY
            p.font.name = "Consolas"

        y += 0.42

    # Quick examples section
    y += 0.2
    box = slide.shapes.add_textbox(Inches(0.3), Inches(y), Inches(9.4), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Quick Examples"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    examples = [
        "mbo /path/to/data.tiff                    # View TIFF in GUI",
        "mbo convert raw/ out.zarr -e .zarr       # Convert to Zarr",
        "mbo convert data.tiff out/ --fix-phase   # Fix bidirectional scan phase",
        "mbo view data/ --roi 0,1                 # View specific ROIs",
    ]

    y += 0.4
    for ex in examples:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9.2), Inches(0.32))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = ex
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.name = "Consolas"
        y += 0.32

    return slide


def add_gui_slide(prs, title, features, image_placeholder_text):
    """Add GUI slide with image placeholder."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_YELLOW

    # Image placeholder (left side)
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(0.8),
        Inches(5.5), Inches(4.2)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = CODE_BG
    placeholder.line.color.rgb = ACCENT_CYAN

    # Placeholder text
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5.1), Inches(0.8))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"[Insert screenshot: {image_placeholder_text}]"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    # Features list (right side)
    features_box = slide.shapes.add_textbox(Inches(6.0), Inches(0.8), Inches(3.7), Inches(4.2))
    tf = features_box.text_frame
    tf.word_wrap = True

    for i, feature in enumerate(features):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(feature, tuple):
            p.text = f"▸ {feature[0]}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT_CYAN

            p2 = tf.add_paragraph()
            p2.text = f"   {feature[1]}"
            p2.font.size = Pt(11)
            p2.font.color.rgb = LIGHT_GRAY
        else:
            p.text = f"▸ {feature}"
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE

    return slide


def add_formats_slide(prs):
    """Add supported formats slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Supported Formats"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_YELLOW

    # Input formats
    box = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(4.5), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Input Formats"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    inputs = [
        (".tif, .tiff", "BigTIFF, OME-TIFF, ScanImage"),
        (".zarr", "Zarr v3, OME-NGFF"),
        (".h5, .hdf5", "HDF5 datasets"),
        (".bin", "Suite2p binary + ops.npy"),
        (".npy", "NumPy arrays"),
        ("In-memory", "NumPy/Dask arrays"),
    ]

    y = 1.25
    for ext, desc in inputs:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(1.3), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = ext
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = "Consolas"

        box = slide.shapes.add_textbox(Inches(1.9), Inches(y), Inches(2.8), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        y += 0.32

    # Output formats
    box = slide.shapes.add_textbox(Inches(5.2), Inches(0.8), Inches(4.5), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Output Formats"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    outputs = [
        (".tiff", "BigTIFF (streaming write)"),
        (".zarr", "Zarr v3 with OME metadata"),
        (".h5", "HDF5 with chunking"),
        (".bin", "Suite2p binary format"),
        (".npy", "NumPy array"),
        (".mp4", "Video export"),
    ]

    y = 1.25
    for ext, desc in outputs:
        box = slide.shapes.add_textbox(Inches(5.4), Inches(y), Inches(1.0), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = ext
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = "Consolas"

        box = slide.shapes.add_textbox(Inches(6.5), Inches(y), Inches(3.2), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        y += 0.32

    # Array types section
    box = slide.shapes.add_textbox(Inches(0.3), Inches(3.6), Inches(9.4), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Lazy Array Types (returned by imread)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    arrays = [
        "MboRawArray - Raw ScanImage multi-ROI data with metadata",
        "TiffArray - Memory-mapped TIFF access",
        "ZarrArray - Chunked cloud-ready arrays",
        "Suite2pArray - Suite2p binary with ops integration",
        "H5Array - HDF5 dataset wrapper",
    ]

    y = 4.05
    for arr in arrays:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9.0), Inches(0.28))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {arr}"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        y += 0.28

    return slide


def create_cheatsheet():
    """Create the complete cheat sheet."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # Slide 1: Title
    add_title_slide(prs, "mbo_utilities", "Cheat Sheet - Python API, CLI & GUI")

    # Slide 2: Core I/O Functions
    add_section_slide(
        prs,
        "Python API: Core I/O",
        [
            ("imread(path)", "Lazy-load any supported format"),
            ("imwrite(arr, path, ext)", "Stream-write to disk"),
            ("get_metadata(path)", "Extract file metadata dict"),
            ("get_voxel_size(path)", "Get physical dimensions (µm)"),
            ("get_files(path)", "Discover files with filtering"),
        ],
        [
            ("Load data", "from mbo_utilities import imread\narr = imread('/path/to/data.tiff')\narr = imread('/path/to/raw/')  # ScanImage"),
            ("Write data", "from mbo_utilities import imwrite\nimwrite(arr, 'out.zarr', ext='.zarr')\nimwrite(arr, 'out/', planes=[0,1,2])"),
            ("Get info", "from mbo_utilities import get_metadata\nmeta = get_metadata('/path/to/data')\nprint(meta['nframes'], meta['shape'])"),
        ]
    )

    # Slide 3: Utility Functions
    add_section_slide(
        prs,
        "Python API: Utilities & Visualization",
        [
            ("save_mp4(fname, images)", "Export video from 3D array"),
            ("save_png(fname, data)", "Save image via matplotlib"),
            ("norm_minmax(images)", "Normalize to 0-1 range"),
            ("smooth_data(data, window)", "Temporal smoothing"),
            ("subsample_array(arr, factor)", "Downsample spatially/temporally"),
            ("files_to_dask(files)", "Build Dask array from files"),
            ("expand_paths(paths)", "Expand wildcards/lists"),
        ],
        [
            ("Video export", "from mbo_utilities import save_mp4\nsave_mp4('movie.mp4', arr[:500],\n         framerate=30,\n         temporal_avg=5)"),
            ("Dask arrays", "from mbo_utilities import files_to_dask\ndarr = files_to_dask(tiff_files,\n                     chunk_t=250)"),
            ("Normalize", "from mbo_utilities import norm_minmax\nnormed = norm_minmax(images)"),
        ]
    )

    # Slide 4: CLI Commands
    add_cli_slide(prs)

    # Slide 5: Supported Formats
    add_formats_slide(prs)

    # Slide 6: GUI - Preview Widget
    add_gui_slide(
        prs,
        "GUI: Preview & Visualization",
        [
            ("Image Viewer", "FastPlotLib 2D/3D rendering"),
            ("Frame Navigation", "Time slider with playback"),
            ("Z-Plane Slider", "Navigate through planes"),
            ("Window Functions", "mean, max, std, mean-sub"),
            ("Scan-Phase Correction", "Fix bidirectional artifacts"),
            ("Contrast Controls", "V-Min/V-Max adjustment"),
            ("Summary Stats", "Per-plane mean, std, SNR"),
        ],
        "docs/_images/GUI_Slide1.png"
    )

    # Slide 7: GUI - Processing
    add_gui_slide(
        prs,
        "GUI: Processing & Export",
        [
            ("Spatial Crop", "Select ROI for processing"),
            ("Suite2p Pipeline", "Integrated registration & detection"),
            ("Registration Settings", "Rigid/non-rigid, 1P mode"),
            ("Save As Dialog", "Export to .tiff/.zarr/.h5/.bin"),
            ("Multi-ROI Support", "Process ROIs separately"),
            ("Suite3D Registration", "Axial z-plane alignment"),
            ("Phase Correction", "Bidirectional scan fix"),
        ],
        "docs/_images/GUI_Slide2.png"
    )

    # Slide 8: GUI - Diagnostics
    add_gui_slide(
        prs,
        "GUI: ROI Diagnostics (Suite2p Results)",
        [
            ("dF/F Traces", "Adjustable baseline method"),
            ("Quality Metrics", "SNR, skewness, activity"),
            ("Filter Histograms", "Interactive threshold sliders"),
            ("Auto-save", "Syncs iscell.npy to disk"),
            ("File Watching", "Detects external changes"),
            ("Suite2p Sync", "Bi-directional with GUI"),
            ("ROI Statistics", "Detailed per-ROI info"),
        ],
        "DiagnosticsWidget"
    )

    # Save
    out_path = Path(__file__).parent.parent / "cheat_sheet.pptx"
    prs.save(out_path)
    print(f"Created: {out_path}")
    return out_path


if __name__ == "__main__":
    create_cheatsheet()
